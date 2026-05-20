"""
Darwish Interserve – Executive Management Presentation Builder
Generates: DarwishInterserve_Executive_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x3A, 0x6B)
NAVY_DARK   = RGBColor(0x12, 0x26, 0x49)
NAVY_LIGHT  = RGBColor(0x2A, 0x52, 0x98)
GOLD        = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_LIGHT  = RGBColor(0xE8, 0xC9, 0x6A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xF5, 0xF6, 0xFA)
TEXT_SEC    = RGBColor(0x6B, 0x72, 0x80)
TEXT_LIGHT  = RGBColor(0x9C, 0xA3, 0xAF)
SUCCESS     = RGBColor(0x10, 0xB9, 0x81)
ERROR       = RGBColor(0xEF, 0x44, 0x44)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
TEAL        = RGBColor(0x14, 0xB8, 0xA6)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)

# ── Slide size: Widescreen 16:9 ────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ─── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_para(tf, text, size=13, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p

def gold_bar(slide, x, y, w=0.4, h=0.045):
    add_rect(slide, x, y, w, h, GOLD)

def section_eyebrow(slide, text, x, y):
    add_text(slide, text, x, y, 5, 0.3, size=9, bold=True, color=GOLD,
             font_name="Calibri")

def stat_block(slide, value, label, x, y, val_color=GOLD, bg=None):
    if bg:
        add_rect(slide, x, y, 1.9, 1.05, bg)
    add_text(slide, value, x, y + 0.08, 1.9, 0.5, size=28, bold=True,
             color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.58, 1.9, 0.35, size=10, bold=False,
             color=TEXT_SEC, align=PP_ALIGN.CENTER)

def card_rect(slide, x, y, w, h, fill=WHITE, radius=True):
    s = slide.shapes.add_shape(5 if radius else 1,   # rounded rect = 5
                               Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    s.line.width = Pt(0.5)
    if radius:
        s.adjustments[0] = 0.05
    return s

def divider(slide, x, y, w, color=RGBColor(0xE5,0xE7,0xEB)):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

# Full dark navy background
add_rect(sl, 0, 0, 13.33, 7.5, NAVY_DARK)

# Gold accent left stripe
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

# Large decorative circle (top-right)
c = sl.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x2A,0x52,0x98)
c.line.fill.background()
# semi-transparent effect via alpha workaround – just lighten
c.fill.fore_color.rgb = RGBColor(0x22,0x44,0x7A)

# Gold accent circle (bottom-left)
c2 = sl.shapes.add_shape(9, Inches(-0.8), Inches(5.5), Inches(3.5), Inches(3.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xC9,0xA8,0x4C)
c2.line.fill.background()

# DI Logo box
add_rect(sl, 0.6, 1.3, 0.9, 0.9, WHITE)
add_text(sl, "DI", 0.6, 1.35, 0.9, 0.7, size=28, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER)
add_rect(sl, 0.6, 2.05, 0.9, 0.12, GOLD)   # gold bar under logo

# Company wordmark
add_text(sl, "DARWISH", 1.65, 1.32, 4, 0.45, size=22, bold=True, color=WHITE,
         font_name="Calibri")
add_text(sl, "INTERSERVE", 1.65, 1.72, 4, 0.38, size=16, bold=True,
         color=GOLD_LIGHT, font_name="Calibri")
add_text(sl, "Facility Management", 1.65, 2.05, 4, 0.3, size=11,
         color=RGBColor(0xCC,0xCC,0xCC), font_name="Calibri")

# Gold divider
add_rect(sl, 0.6, 2.65, 2.2, 0.055, GOLD)

# Title
add_text(sl, "Mobile Application", 0.6, 2.88, 9, 0.7, size=38, bold=True,
         color=WHITE, font_name="Calibri")
add_text(sl, "Executive Presentation", 0.6, 3.5, 9, 0.55, size=28, bold=False,
         color=GOLD_LIGHT, font_name="Calibri")

# Subtitle
add_text(sl,
         "Integrated Facility Management Platform  ·  React Native (Expo)  ·  Qatar 2026",
         0.6, 4.2, 10, 0.4, size=13, color=RGBColor(0xAA,0xBB,0xCC),
         font_name="Calibri")

# Stat pills bottom
stats = [("3,500+","Employees"),("100+","Active Projects"),("14+","Years in Qatar"),("235","Active Sites")]
for i,(v,l) in enumerate(stats):
    bx = 0.6 + i*3.1
    add_rect(sl, bx, 5.5, 2.75, 1.1, RGBColor(0x1F,0x42,0x75))
    add_text(sl, v, bx, 5.55, 2.75, 0.55, size=24, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(sl, l, bx, 6.08, 2.75, 0.35, size=10, color=RGBColor(0xCC,0xDD,0xEE),
             align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(sl, 0, 7.15, 13.33, 0.35, NAVY)
add_text(sl, "CONFIDENTIAL — FOR MANAGEMENT USE ONLY", 0, 7.18, 13.33, 0.3,
         size=9, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.5, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "AGENDA", 0.45, 0.12, 6, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "What We'll Cover Today", 0.45, 0.42, 10, 0.7, size=30,
         bold=True, color=WHITE)

items = [
    ("01", "Company Overview", "Who we are, our vision, mission & certifications"),
    ("02", "App Purpose & User Journey", "Guest & Employee authentication flows"),
    ("03", "Home Dashboard & Sub-Tabs", "Website · Management Team · Heads · Employees"),
    ("04", "Services Portfolio", "Hard, Soft & Managed FM services"),
    ("05", "Projects Portfolio", "Active engagements across sectors"),
    ("06", "News & Events", "Internal communications platform"),
    ("07", "Leadership Hierarchy", "Org chart from CEO to Department Heads"),
    ("08", "Employee Management", "Searchable staff directory with status"),
    ("09", "Profile & Company About", "Vision, Mission, Values, ISO certs, Contact"),
    ("10", "Technology & Next Steps", "Tech stack, roadmap & call to action"),
]

cols = [items[:5], items[5:]]
for ci, col in enumerate(cols):
    cx = 0.45 + ci * 6.5
    for ri, (num, title, sub) in enumerate(col):
        ry = 1.75 + ri * 1.02
        card_rect(sl, cx, ry, 6.1, 0.88, WHITE)
        add_rect(sl, cx, ry, 0.52, 0.88, NAVY)
        add_text(sl, num, cx, ry + 0.18, 0.52, 0.5, size=17, bold=True,
                 color=GOLD, align=PP_ALIGN.CENTER)
        add_text(sl, title, cx + 0.6, ry + 0.07, 5.3, 0.38, size=13, bold=True,
                 color=NAVY_DARK)
        add_text(sl, sub,   cx + 0.6, ry + 0.46, 5.3, 0.3, size=10,
                 color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — COMPANY OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "COMPANY OVERVIEW", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Who We Are", 0.45, 0.42, 10, 0.7, size=30, bold=True, color=WHITE)

# About card
card_rect(sl, 0.45, 1.65, 5.9, 2.05, WHITE)
add_text(sl, "WHO WE ARE", 0.65, 1.78, 5.5, 0.3, size=9, bold=True, color=GOLD)
add_text(sl, "A 100% Qatari-Owned Integrated FM Company",
         0.65, 2.05, 5.5, 0.45, size=14, bold=True, color=NAVY_DARK)
add_text(sl,
         "Darwish Interserve Facility Management (DIFM) is a joint venture "
         "between Darwish Holdings and Interserve, delivering world-class "
         "integrated facility management services across Qatar since 2010. "
         "We serve government, aviation, education, and commercial sectors.",
         0.65, 2.5, 5.6, 0.95, size=11, color=TEXT_SEC)

# Vision / Mission cards
card_rect(sl, 0.45, 3.85, 2.83, 1.65, RGBColor(0xEE,0xF2,0xFF))
add_text(sl, "VISION", 0.65, 3.95, 2.5, 0.3, size=9, bold=True, color=NAVY)
add_text(sl,
         "To be the leading FM provider in Qatar, delivering excellence through "
         "innovation, quality, and customer satisfaction.",
         0.65, 4.22, 2.5, 0.95, size=10.5, color=NAVY_DARK)

card_rect(sl, 3.45, 3.85, 2.83, 1.65, RGBColor(0xFF,0xF8,0xE7))
add_text(sl, "MISSION", 3.65, 3.95, 2.5, 0.3, size=9, bold=True, color=GOLD)
add_text(sl,
         "To provide reliable and efficient FM services while maintaining the "
         "highest standards of safety, quality, and sustainability.",
         3.65, 4.22, 2.5, 0.95, size=10.5, color=NAVY_DARK)

# Right — stats
stats_r = [
    ("3,500+","Skilled Employees", BLUE),
    ("100+","Active Projects",     GOLD),
    ("14+","Years of Excellence",  SUCCESS),
    ("235","Active Sites",         NAVY_LIGHT),
]
for i,(v,l,c) in enumerate(stats_r):
    bx = 6.7 + (i%2)*3.2
    by = 1.65 + (i//2)*2.0
    card_rect(sl, bx, by, 2.95, 1.72, WHITE)
    add_rect(sl, bx, by, 2.95, 0.18, c)
    add_text(sl, v, bx, by + 0.28, 2.95, 0.65, size=34, bold=True,
             color=c, align=PP_ALIGN.CENTER)
    add_text(sl, l, bx, by + 0.98, 2.95, 0.4, size=11, color=TEXT_SEC,
             align=PP_ALIGN.CENTER)

# ISO certs row
certs = [("ISO 9001","Quality Mgmt",BLUE),("ISO 14001","Environment",SUCCESS),
         ("ISO 45001","H&S",ERROR),("ISO 41001","FM Standards",PURPLE)]
for i,(code,lbl,c) in enumerate(certs):
    bx = 0.45 + i * 3.22
    card_rect(sl, bx, 5.7, 3.0, 0.92, WHITE)
    add_rect(sl, bx, 5.7, 3.0, 0.12, c)
    add_text(sl, code, bx, 5.8,  3.0, 0.38, size=15, bold=True, color=c,
             align=PP_ALIGN.CENTER)
    add_text(sl, lbl,  bx, 6.16, 3.0, 0.3,  size=10, color=TEXT_SEC,
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — APP OVERVIEW & USER JOURNEY
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "APP OVERVIEW & USER JOURNEY", 0.45, 0.12, 9, 0.32, size=10,
         bold=True, color=GOLD)
add_text(sl, "Dual-Access Authentication Platform", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

# Left column — purpose
card_rect(sl, 0.45, 1.65, 5.5, 1.65, WHITE)
add_text(sl, "PURPOSE", 0.65, 1.75, 5, 0.3, size=9, bold=True, color=GOLD)
add_text(sl, "What the App Solves", 0.65, 2.0, 5, 0.38, size=15, bold=True,
         color=NAVY_DARK)
add_text(sl,
         "A single, secure mobile gateway for both internal employees and "
         "external guests to access facility management information, "
         "project updates, events, and staff directories.",
         0.65, 2.38, 5.1, 0.75, size=11, color=TEXT_SEC)

# User types
for i,(icon,title,desc,c) in enumerate([
    ("👤","Guest Access","Gmail OTP · No account required · View-only",
     RGBColor(0xFF,0xF8,0xE7)),
    ("🏢","Employee Access","SMS OTP · Full dashboard · Work orders · Profile",
     RGBColor(0xEE,0xF2,0xFF)),
]):
    card_rect(sl, 0.45, 3.5 + i*1.65, 5.5, 1.42, c)
    add_text(sl, icon + "  " + title, 0.65, 3.62 + i*1.65, 5, 0.38,
             size=15, bold=True, color=NAVY_DARK)
    add_text(sl, desc, 0.65, 4.0 + i*1.65, 5, 0.7, size=11, color=TEXT_SEC)

# Flow arrows — right side
steps = [
    ("01","Welcome Screen","Brand landing with Guest & Employee entry points",NAVY),
    ("02","Authentication","Enter credentials · OTP verification (Gmail / SMS)",NAVY_LIGHT),
    ("03","OTP Verification","6-digit code · Auto-advance · Countdown resend",GOLD),
    ("04","Main Dashboard","Home · Events · Projects · Profile tabs",SUCCESS),
    ("05","Content & Actions","Sub-tabs · Search · Project detail · Notifications",BLUE),
]
for i,(num,title,desc,c) in enumerate(steps):
    by = 1.62 + i * 1.13
    card_rect(sl, 6.25, by, 6.75, 0.98, WHITE)
    add_rect(sl, 6.25, by, 0.52, 0.98, c)
    add_text(sl, num, 6.25, by + 0.22, 0.52, 0.5, size=16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 6.85, by + 0.06, 6.0, 0.38, size=13, bold=True,
             color=NAVY_DARK)
    add_text(sl, desc,  6.85, by + 0.44, 6.0, 0.42, size=10, color=TEXT_SEC)
    if i < 4:
        add_rect(sl, 6.47, by + 1.0, 0.07, 0.13, GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — HOME DASHBOARD & SUB-TABS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "HOME DASHBOARD", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "4 Interactive Sub-Tabs", 0.45, 0.42, 10, 0.7, size=30,
         bold=True, color=WHITE)

tabs = [
    ("01","Website",
     ["Company hero & welcome banner","About Us — DIFM overview",
      "Our Services (Hard, Soft, Managed)","By The Numbers stats grid",
      "Client logos (HIA, QU, Katara…)","CTA: Explore Projects · Contact Us"],
     NAVY),
    ("02","Management Team",
     ["Visual org-chart hierarchy","CEO at the apex",
      "COO → HR Director & Ops Director","CFO  ·  CTO direct reports",
      "Colour-coded reporting lines","Role badges on each card"],
     NAVY_LIGHT),
    ("03","Management Heads",
     ["8 department head detail cards","Reporting line clearly shown",
      "Team size & active site count","Key responsibilities per head",
      "Ordered by reporting structure","Colour-coded by department"],
     GOLD),
    ("04","Employee List",
     ["Searchable directory of all staff","Real-time filter by name / dept / ID",
      "Employee ID (DI-2024-XXX) shown","Active / On Leave status badges",
      "Department grouping","Live count updates on search"],
     SUCCESS),
]

for i,(num,title,bullets,c) in enumerate(tabs):
    bx = 0.45 + i * 3.21
    # Header card
    add_rect(sl, bx, 1.65, 3.0, 0.65, c)
    add_text(sl, num,   bx,       1.68, 0.55, 0.5, size=14, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(sl, title, bx + 0.55, 1.7, 2.35, 0.5, size=15, bold=True,
             color=WHITE)
    # Bullets
    card_rect(sl, bx, 2.3, 3.0, 4.85, WHITE)
    for j, b in enumerate(bullets):
        by = 2.45 + j * 0.72
        add_rect(sl, bx + 0.2, by + 0.13, 0.07, 0.07, c)
        add_text(sl, b, bx + 0.38, by, 2.48, 0.6, size=11, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — SERVICES PORTFOLIO
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "SERVICES PORTFOLIO", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Integrated Facility Management Services", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

services = [
    ("Hard Services", BLUE,
     ["Mechanical & Electrical","HVAC Systems","Civil & Structural",
      "Plumbing & Fire Safety","Preventive Maintenance"]),
    ("Soft Services", SUCCESS,
     ["Cleaning & Housekeeping","Landscaping & Horticulture",
      "Pest Control","Security & Guarding","Support Services"]),
    ("Managed Services", PURPLE,
     ["Energy Management","CAFM / CMMS Platform",
      "Compliance & Auditing","Consulting & Advisory",
      "International Standards"]),
    ("MEP Engineering", GOLD,
     ["Mechanical Systems","Electrical Infrastructure",
      "Plumbing Networks","BMS Integration","Asset Management"]),
    ("Cleaning", TEAL,
     ["Daily Housekeeping","Deep Cleaning Cycles",
      "Specialised Cleaning","Waste Management","Hygiene Compliance"]),
    ("Security", ERROR,
     ["Access Control","CCTV & Monitoring",
      "Manned Guarding","Incident Response","Emergency Protocols"]),
]

for i,( title, c, bullets) in enumerate(services):
    bx = 0.45 + (i%3) * 4.3
    by = 1.65 + (i//3) * 2.75
    card_rect(sl, bx, by, 4.05, 2.55, WHITE)
    add_rect(sl, bx, by, 4.05, 0.45, c)
    add_text(sl, title, bx + 0.15, by + 0.06, 3.7, 0.35, size=14, bold=True,
             color=WHITE)
    for j,b in enumerate(bullets):
        add_rect(sl, bx + 0.18, by + 0.62 + j*0.37, 0.06, 0.06, c)
        add_text(sl, b, bx + 0.34, by + 0.52 + j*0.37, 3.55, 0.34,
                 size=11, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — PROJECTS PORTFOLIO
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "PROJECTS PORTFOLIO", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Active Engagements Across Sectors", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

# KPIs
kpis = [("100+","Active Projects",NAVY),("235","Total Sites",BLUE),
        ("3,500+","Workforce",SUCCESS),("14+","Years",GOLD)]
for i,(v,l,c) in enumerate(kpis):
    bx = 0.45 + i * 3.21
    card_rect(sl, bx, 1.62, 3.0, 1.1, WHITE)
    add_rect(sl, bx, 1.62, 3.0, 0.12, c)
    add_text(sl, v, bx, 1.72, 3.0, 0.6, size=30, bold=True, color=c,
             align=PP_ALIGN.CENTER)
    add_text(sl, l, bx, 2.3, 3.0, 0.3, size=10, color=TEXT_SEC,
             align=PP_ALIGN.CENTER)

projects = [
    ("01","Integrated FM — Education Sector",
     "215 sites · 1,200+ staff","Qatar Public Schools · MoE","92%",SUCCESS,
     "Qatar's largest FM contract covering schools nationwide with full soft "
     "and hard services, preventive maintenance, and compliance."),
    ("02","Aviation Facility Maintenance",
     "12 sites · 240 staff","Hamad International Airport","78%",BLUE,
     "Specialised HVAC, MEP, and fire-safety systems maintenance for one "
     "of the world's top-rated international airports."),
    ("03","Commercial Property Management",
     "8 sites · 95 staff","Corporate Clients","35%",PURPLE,
     "End-to-end facility management ensuring operational efficiency and "
     "tenant satisfaction across premium commercial properties."),
]

for i,(num,title,meta,client,pct,c,desc) in enumerate(projects):
    bx = 0.45 + i * 4.3
    card_rect(sl, bx, 2.9, 4.05, 4.35, WHITE)
    add_rect(sl, bx, 2.9, 4.05, 0.55, c)
    add_text(sl, num,    bx,       2.93, 0.6, 0.48, size=16, bold=True,
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(sl, title,  bx + 0.62, 2.97, 3.3, 0.44, size=12, bold=True,
             color=WHITE)
    add_text(sl, "CLIENT: " + client, bx + 0.15, 3.52, 3.7, 0.3, size=10,
             bold=True, color=c)
    add_text(sl, meta,   bx + 0.15, 3.8,  3.7, 0.3, size=10, color=TEXT_SEC)
    add_text(sl, desc,   bx + 0.15, 4.1,  3.7, 0.85, size=10.5, color=TEXT_SEC)
    # progress bar bg
    add_rect(sl, bx + 0.15, 5.08, 3.7, 0.14, RGBColor(0xE5,0xE7,0xEB))
    fill_w = 3.7 * int(pct.replace('%','')) / 100
    add_rect(sl, bx + 0.15, 5.08, fill_w, 0.14, c)
    add_text(sl, "Progress  " + pct, bx + 0.15, 4.9, 3.7, 0.3, size=10,
             bold=True, color=c)
    add_text(sl, "View in App →", bx + 0.15, 5.36, 3.7, 0.3, size=10,
             bold=True, color=c)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — NEWS & EVENTS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "NEWS & EVENTS", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Internal Communications Platform", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

# Featured event card
card_rect(sl, 0.45, 1.65, 7.5, 3.2, NAVY)
add_text(sl, "★  FEATURED EVENT", 0.65, 1.78, 7, 0.35, size=10, bold=True,
         color=GOLD)
add_text(sl, "Mandatory HSE Safety Training",
         0.65, 2.12, 7, 0.65, size=22, bold=True, color=WHITE)
add_text(sl, "📅  May 10, 2026  ·  09:00 – 12:00",
         0.65, 2.77, 7, 0.35, size=12, color=GOLD_LIGHT)
add_text(sl, "📍  Main Office, Doha  ·  +142 attending",
         0.65, 3.1,  7, 0.35, size=12, color=GOLD_LIGHT)
add_text(sl, "Mandatory attendance for all site supervisors and HSE officers. "
         "Covers updated Qatar Civil Defence requirements and incident response protocols.",
         0.65, 3.48, 7, 0.6, size=11.5, color=RGBColor(0xCC,0xDD,0xEE))
add_rect(sl, 0.65, 4.35, 7.1, 0.01, RGBColor(0x33,0x44,0x66))
add_text(sl, "RSVP →", 0.65, 4.42, 1.5, 0.3, size=11, bold=True, color=GOLD)

# Filter categories
for i,(lbl,c) in enumerate([("All","active"),("Training",""),
                              ("HSE",""),("Company",""),("Holidays","")]):
    bx = 0.45 + i * 1.53
    add_rect(sl, bx, 5.08, 1.38, 0.38,
             NAVY if lbl=="All" else WHITE)
    add_text(sl, lbl, bx, 5.1, 1.38, 0.34, size=11,
             bold=(lbl=="All"),
             color=WHITE if lbl=="All" else TEXT_SEC,
             align=PP_ALIGN.CENTER)

# Event list
events = [
    ("12\nMAY","Training","New Work Order System Walkthrough",
     "14:00–15:30  ·  Online – MS Teams",BLUE),
    ("18\nMAY","Company","Q2 Town Hall Meeting",
     "10:00–11:30  ·  DIFM Auditorium",PURPLE),
    ("27\nMAY","Holiday","Eid Al-Adha Holiday",
     "All Day  ·  Company-wide",SUCCESS),
]
for i,(date,tag,title,meta,c) in enumerate(events):
    bx = 8.15 + 0
    by = 1.65 + i * 1.88
    card_rect(sl, 0.45, by + 5.62 - 5.62, 12.55 - 8.15 + 0.18, 1.65, WHITE)

# Redo event list on right
for i,(date,tag,title,meta,c) in enumerate(events):
    by = 1.65 + i * 1.9
    card_rect(sl, 8.25, by, 4.75, 1.65, WHITE)
    add_rect(sl, 8.25, by, 0.75, 1.65, RGBColor(220, 230, 245))
    # date
    add_text(sl, date.split('\n')[0], 8.25, by + 0.28, 0.75, 0.45,
             size=18, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(sl, date.split('\n')[1], 8.25, by + 0.7, 0.75, 0.3,
             size=9, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_rect(sl, 8.28, by + 0.12, 0.68, 0.08, RGBColor(200, 215, 235))
    add_text(sl, "● " + tag, 9.1, by + 0.1, 3.7, 0.3, size=10, bold=True, color=c)
    add_text(sl, title, 9.1, by + 0.4, 3.7, 0.45, size=12, bold=True,
             color=NAVY_DARK)
    add_text(sl, meta,  9.1, by + 0.85, 3.7, 0.3, size=10, color=TEXT_SEC)

# Stats strip under featured
for i,(v,l,c) in enumerate([("5","Upcoming Events",GOLD),
                              ("2","New This Week",BLUE),
                              ("142","RSVPs Confirmed",SUCCESS)]):
    bx = 0.45 + i * 2.55
    card_rect(sl, bx, 5.62, 2.35, 1.0, WHITE)
    add_rect(sl, bx, 5.62, 2.35, 0.1, c)
    add_text(sl, v, bx, 5.7, 2.35, 0.5, size=24, bold=True, color=c,
             align=PP_ALIGN.CENTER)
    add_text(sl, l, bx, 6.18, 2.35, 0.3, size=9, color=TEXT_SEC,
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — LEADERSHIP HIERARCHY
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "LEADERSHIP HIERARCHY", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Organisational Structure", 0.45, 0.42, 11, 0.7, size=30,
         bold=True, color=WHITE)

def org_card(slide, x, y, w, initials, name, title, c, badge, badge_c, badge_bg):
    card_rect(slide, x, y, w, 1.05, WHITE)
    add_rect(slide, x, y, w, 0.1, c)
    # avatar circle
    av = slide.shapes.add_shape(9, Inches(x+0.12), Inches(y+0.2),
                                 Inches(0.6), Inches(0.6))
    av.fill.solid(); av.fill.fore_color.rgb = c
    av.line.fill.background()
    add_text(slide, initials, x+0.12, y+0.28, 0.6, 0.42, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name,  x+0.82, y+0.12, w-0.98, 0.38, size=11, bold=True,
             color=NAVY_DARK)
    add_text(slide, title, x+0.82, y+0.5,  w-0.98, 0.38, size=9,  color=TEXT_SEC)
    badge_s = slide.shapes.add_shape(5, Inches(x+w-0.72), Inches(y+0.3),
                                      Inches(0.62), Inches(0.3))
    badge_s.fill.solid(); badge_s.fill.fore_color.rgb = badge_bg
    badge_s.line.fill.background()
    if badge_s.adjustments:
        badge_s.adjustments[0] = 0.5
    add_text(slide, badge, x+w-0.72, y+0.3, 0.62, 0.3, size=9, bold=True,
             color=badge_c, align=PP_ALIGN.CENTER)

def v_line(slide, x, y, h):
    add_rect(slide, x, y, 0.025, h, RGBColor(0xD1,0xD5,0xDB))

def h_line(slide, x, y, w):
    add_rect(slide, x, y, w, 0.025, RGBColor(0xD1,0xD5,0xDB))

# CEO (centred)
org_card(sl, 4.9, 1.62, 3.5, "KD",
         "Khalid Al-Darwish", "Chief Executive Officer",
         NAVY, "CEO", NAVY, RGBColor(0xEE,0xF2,0xFF))

# Vertical from CEO
v_line(sl, 6.64, 2.67, 0.42)

# Horizontal spanning COO-CFO-CTO
h_line(sl, 2.0, 3.09, 9.3)

# COO
v_line(sl, 2.75, 3.09, 0.35)
org_card(sl, 0.45, 3.44, 4.6, "JR",
         "James Richardson", "Chief Operating Officer",
         SUCCESS, "COO", RGBColor(0x06,0x5f,0x46), RGBColor(0xD1,0xFA,0xE5))

# CFO
v_line(sl, 6.64, 3.09, 0.35)
org_card(sl, 4.9, 3.44, 3.5, "AH",
         "Ahmed Hassan", "Chief Financial Officer",
         GOLD, "CFO", RGBColor(0x92,0x40,0x0e), RGBColor(0xFE,0xF3,0xC7))

# CTO
v_line(sl, 10.99, 3.09, 0.35)
org_card(sl, 9.4, 3.44, 3.5, "MQ",
         "Mohammed Al-Qahtani","Chief Technology Officer",
         PURPLE, "CTO", RGBColor(0x4c,0x1d,0x95), RGBColor(0xED,0xE9,0xFE))

# Under COO: H-line → HR & Ops
v_line(sl, 2.75, 4.49, 0.36)
h_line(sl, 0.75, 4.85, 3.9)
v_line(sl, 0.75, 4.85, 0.32)
v_line(sl, 4.65, 4.85, 0.32)

org_card(sl, 0.45, 5.17, 4.0, "SM",
         "Sara Al-Mansoori","HR Director",
         RGBColor(0xec,0x48,0x99), "HR Dir",
         RGBColor(0x9d,0x17,0x4d), RGBColor(0xFC,0xE7,0xF3))
org_card(sl, 4.65, 5.17, 4.0, "DT",
         "David Thompson","Operations Director",
         BLUE, "Ops Dir", RGBColor(0x1e,0x40,0xaf), RGBColor(0xDB,0xEA,0xFE))

# Under CFO: Rania & Aisha
v_line(sl, 6.64, 4.49, 0.36)
h_line(sl, 4.9, 4.85, 3.5)

# Under CTO: Omar
v_line(sl, 10.99, 4.49, 0.68)
org_card(sl, 9.4, 5.17, 3.5, "OS",
         "Omar Siddiqui","Head of IT & Systems",
         PURPLE, "IT Head",
         RGBColor(0x4c,0x1d,0x95), RGBColor(0xED,0xE9,0xFE))

# CFO sub
v_line(sl, 5.1,  4.85, 0.32)
v_line(sl, 8.2,  4.85, 0.32)
org_card(sl, 4.9, 5.17, 3.1, "RM",
         "Rania Mahmoud","Head of Finance",
         GOLD, "Finance",
         RGBColor(0x92,0x40,0x0e), RGBColor(0xFE,0xF3,0xC7))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — MANAGEMENT HEADS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "MANAGEMENT HEADS", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Department Heads & Team Details", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

heads = [
    ("IK","Ibrahim Al-Khalidi","Hard Services",NAVY,    "COO","320 staff · 45+ sites","MEP · Civil · HVAC · PM"),
    ("NF","Nadia Al-Farsi",   "Soft Services",SUCCESS, "COO","280 staff · 38+ sites","Cleaning · Landscaping · Security"),
    ("FK","Faisal Al-Kuwari", "Projects & Eng.",TEAL,  "Ops Dir","95 engineers · 100+ projects","Fit-Out · Refurb · Delivery"),
    ("RM","Rania Mahmoud",    "Finance",       GOLD,   "CFO","28 staff · QAR 500M+","Budgeting · Payroll · Reporting"),
    ("AT","Aisha Al-Thani",   "Procurement",   RGBColor(0xec,0x48,0x99), "CFO","18 staff · 200+ vendors","Contracts · Sourcing"),
    ("OS","Omar Siddiqui",    "IT & Systems",  PURPLE, "CTO","24 staff","CAFM · CMMS · Mobile Platform"),
    ("PC","Peter Collins",    "HSE & Safety",  ERROR,  "COO","15 officers","Audits · ISO 45001 · Incidents"),
    ("LB","Linda Brown",      "Client Relations",BLUE, "Ops Dir","12 staff · 50+ accounts","SLA · Client Satisfaction"),
]

for i,(init,name,dept,c,reports,team,resp) in enumerate(heads):
    bx = 0.45 + (i%4) * 3.21
    by = 1.65 + (i//4) * 2.7
    card_rect(sl, bx, by, 3.05, 2.5, WHITE)
    add_rect(sl, bx, by, 3.05, 0.1, c)
    # avatar
    av = sl.shapes.add_shape(9, Inches(bx+0.15), Inches(by+0.2),
                              Inches(0.65), Inches(0.65))
    av.fill.solid(); av.fill.fore_color.rgb = c; av.line.fill.background()
    add_text(sl, init, bx+0.15, by+0.3, 0.65, 0.45, size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, name, bx+0.9, by+0.18, 1.95, 0.45, size=11, bold=True,
             color=NAVY_DARK)
    # dept badge
    add_text(sl, dept, bx+0.15, by+1.0, 2.75, 0.3, size=10, bold=True, color=c)
    add_text(sl, "Reports to: " + reports, bx+0.15, by+1.28, 2.75, 0.28,
             size=9, color=TEXT_LIGHT)
    divider(sl, bx+0.15, by+1.56, 2.7)
    add_text(sl, "Team: " + team, bx+0.15, by+1.62, 2.75, 0.3, size=9.5,
             color=TEXT_SEC)
    add_text(sl, resp, bx+0.15, by+1.92, 2.75, 0.38, size=9, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — EMPLOYEE MANAGEMENT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "EMPLOYEE MANAGEMENT", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Searchable Staff Directory", 0.45, 0.42, 11, 0.7, size=30,
         bold=True, color=WHITE)

# Search bar mock
card_rect(sl, 0.45, 1.65, 8.0, 0.6, WHITE)
add_text(sl, "🔍  Search by name, department or ID…", 0.7, 1.74, 7.6, 0.42,
         size=12, color=TEXT_LIGHT)
add_rect(sl, 8.62, 1.65, 1.4, 0.6, NAVY)
add_text(sl, "Search", 8.62, 1.72, 1.4, 0.46, size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

# Count
add_text(sl, "Showing 12 employees", 0.45, 2.35, 4, 0.3, size=10,
         color=TEXT_LIGHT)

emp = [
    ("MM","Mohammed Al-Mansoori","FM Operations","DI-2024-001",True,NAVY),
    ("AR","Ahmed Al-Rashid",     "Hard Services","DI-2024-002",True,SUCCESS),
    ("SK","Sara Al-Kaabi",       "Soft Services","DI-2024-003",True,RGBColor(0xec,0x48,0x99)),
    ("KS","Khalid Al-Suwaidi",   "HSE & Safety", "DI-2024-004",False,ERROR),
    ("FT","Fatima Al-Thani",     "Finance",      "DI-2024-005",True,GOLD),
    ("OM","Omar Al-Mutawa",      "IT & Systems", "DI-2024-006",True,PURPLE),
    ("LB","Layla Binhendi",      "Procurement",  "DI-2024-007",True,TEAL),
    ("RA","Rami Aziz",           "Projects",     "DI-2024-008",False,BLUE),
    ("NA","Noor Al-Ali",         "Soft Services","DI-2024-009",True,SUCCESS),
    ("HK","Hassan Al-Kawari",    "Hard Services","DI-2024-010",True,NAVY),
    ("DA","Dina Al-Ansari",      "Finance",      "DI-2024-011",True,GOLD),
    ("YM","Yousuf Mansoor",      "Operations",   "DI-2024-012",True,NAVY_LIGHT),
]

cols = [emp[:6], emp[6:]]
for ci, column in enumerate(cols):
    cx = 0.45 + ci * 6.45
    for ri, (init,name,dept,eid,active,c) in enumerate(column):
        ry = 2.72 + ri * 0.78
        card_rect(sl, cx, ry, 6.15, 0.7, WHITE)
        av = sl.shapes.add_shape(9, Inches(cx+0.12), Inches(ry+0.08),
                                  Inches(0.52), Inches(0.52))
        av.fill.solid(); av.fill.fore_color.rgb = c; av.line.fill.background()
        add_text(sl, init, cx+0.12, ry+0.13, 0.52, 0.38, size=11, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, name, cx+0.75, ry+0.06, 3.3, 0.35, size=12, bold=True,
                 color=NAVY_DARK)
        add_text(sl, dept, cx+0.75, ry+0.4,  3.3, 0.25, size=9.5, color=TEXT_SEC)
        add_text(sl, eid,  cx+4.15, ry+0.06, 1.8, 0.3, size=9, color=TEXT_LIGHT,
                 align=PP_ALIGN.RIGHT)
        st_c  = SUCCESS if active else GOLD
        st_bg = RGBColor(0xD1,0xFA,0xE5) if active else RGBColor(0xFE,0xF3,0xC7)
        st_lbl= "Active" if active else "On Leave"
        badge_s = sl.shapes.add_shape(5, Inches(cx+4.15), Inches(ry+0.36),
                                       Inches(1.8), Inches(0.24))
        badge_s.fill.solid(); badge_s.fill.fore_color.rgb = st_bg
        badge_s.line.fill.background()
        add_text(sl, st_lbl, cx+4.15, ry+0.36, 1.8, 0.24, size=9, bold=True,
                 color=st_c, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — PROFILE & COMPANY ABOUT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "PROFILE & COMPANY ABOUT", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "User Profile & Company Information", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

# Left — user profile card
card_rect(sl, 0.45, 1.65, 5.5, 5.5, WHITE)
add_rect(sl, 0.45, 1.65, 5.5, 2.1, NAVY_DARK)
av2 = sl.shapes.add_shape(9, Inches(2.38), Inches(1.9), Inches(1.5), Inches(1.5))
av2.fill.solid(); av2.fill.fore_color.rgb = GOLD; av2.line.fill.background()
add_text(sl, "M", 2.38, 1.98, 1.5, 1.0, size=40, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER)
add_text(sl, "Mohammed Al-Mansoori", 0.6, 3.42, 5.1, 0.42, size=15, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "💼  Employee  ·  DI-2024-001", 0.6, 3.82, 5.1, 0.3, size=11,
         color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# Tiles
tiles = [("248","Tasks Done",RGBColor(0xEE,0xF2,0xFF)),
         ("3","Awards",RGBColor(0xFF,0xF8,0xE7)),
         ("1.2y","Tenure",RGBColor(0xD1,0xFA,0xE5))]
for i,(v,l,bg) in enumerate(tiles):
    bx2 = 0.6 + i*1.73
    add_rect(sl, bx2, 4.28, 1.6, 0.88, bg)
    add_text(sl, v, bx2, 4.3, 1.6, 0.48, size=18, bold=True, color=NAVY_DARK,
             align=PP_ALIGN.CENTER)
    add_text(sl, l, bx2, 4.76, 1.6, 0.28, size=9, color=TEXT_SEC,
             align=PP_ALIGN.CENTER)

# Account rows
rows = [("Personal Information","›"),
        ("Security & Password","›"),
        ("Employee Details","›"),
        ("Push Notifications","ON"),
        ("Language","English")]
for i,(lbl,val) in enumerate(rows):
    by2 = 5.28 + i*0.36
    divider(sl, 0.6, by2, 5.2)
    add_text(sl, lbl, 0.65, by2+0.04, 4.3, 0.3, size=11, color=NAVY_DARK)
    add_text(sl, val, 4.9,  by2+0.04, 0.9, 0.3, size=11,
             color=SUCCESS if val=="ON" else TEXT_SEC,
             bold=(val=="ON"), align=PP_ALIGN.RIGHT)

# Right — Company About
card_rect(sl, 6.2, 1.65, 6.75, 2.55, WHITE)
add_text(sl, "COMPANY OVERVIEW", 6.4, 1.75, 6.3, 0.3, size=9, bold=True, color=GOLD)
add_text(sl, "Darwish Interserve is a 100% Qatari-owned integrated facilities "
         "management company delivering world-class services across multiple "
         "sectors including aviation, education, government and commercial.",
         6.4, 2.0, 6.3, 0.85, size=11, color=TEXT_SEC)

# Vision / Mission
card_rect(sl, 6.2, 4.32, 3.22, 1.52, RGBColor(0xEE,0xF2,0xFF))
add_text(sl, "VISION", 6.4, 4.42, 2.8, 0.28, size=9, bold=True, color=NAVY)
add_text(sl, "To be the leading FM provider in Qatar through excellence and innovation.",
         6.4, 4.68, 2.9, 0.9, size=10.5, color=NAVY_DARK)

card_rect(sl, 9.62, 4.32, 3.22, 1.52, RGBColor(0xFF,0xF8,0xE7))
add_text(sl, "MISSION", 9.82, 4.42, 2.8, 0.28, size=9, bold=True, color=GOLD)
add_text(sl, "To deliver reliable FM services at the highest standards of safety and quality.",
         9.82, 4.68, 2.9, 0.9, size=10.5, color=NAVY_DARK)

# Values row
vals = [("Integrity",BLUE),("Excellence",GOLD),("Innovation",PURPLE),("Teamwork",SUCCESS)]
for i,(v,c) in enumerate(vals):
    bx3 = 6.2 + i*1.7
    card_rect(sl, bx3, 5.96, 1.6, 0.7, WHITE)
    add_rect(sl, bx3, 5.96, 1.6, 0.1, c)
    add_text(sl, v, bx3, 6.05, 1.6, 0.42, size=11, bold=True, color=c,
             align=PP_ALIGN.CENTER)

# Certs below values
certs = [("ISO 9001",BLUE),("ISO 14001",SUCCESS),("ISO 45001",ERROR),("ISO 41001",PURPLE)]
for i,(code,c) in enumerate(certs):
    bx4 = 6.2 + i*1.7
    card_rect(sl, bx4, 6.76, 1.6, 0.58, WHITE)
    add_rect(sl, bx4, 6.76, 1.6, 0.08, c)
    add_text(sl, code, bx4, 6.82, 1.6, 0.42, size=11, bold=True, color=c,
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — TECHNOLOGY & ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "TECHNOLOGY & ARCHITECTURE", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Built for Scale, Speed & Security", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

tech = [
    ("React Native","Cross-platform iOS & Android from a single codebase. "
     "Native performance with JavaScript flexibility.",BLUE,"Frontend"),
    ("Expo ~51","Managed workflow with OTA updates, push notifications, "
     "and deep-link support out of the box.",TEAL,"Framework"),
    ("React Navigation v6","Stack + Bottom-Tab navigation with smooth "
     "slide transitions and history management.",SUCCESS,"Navigation"),
    ("OTP Authentication","Gmail (6-digit) for guests · SMS (6-digit) for "
     "employees · Auto-advance inputs · Countdown resend.",GOLD,"Security"),
    ("CAFM / CMMS","Integrated work-order management, asset tracking, and "
     "preventive maintenance scheduling.",PURPLE,"Operations"),
    ("GitHub Pages","Live interactive HTML prototype hosted on GitHub Pages "
     "for stakeholder demos and reviews.",ERROR,"Deployment"),
]

for i,(name,desc,c,tag) in enumerate(tech):
    bx = 0.45 + (i%3)*4.3
    by = 1.65 + (i//3)*2.65
    card_rect(sl, bx, by, 4.05, 2.42, WHITE)
    add_rect(sl, bx, by, 4.05, 0.5, c)
    add_text(sl, tag,  bx+0.15, by+0.06, 2.0, 0.25, size=9, bold=True, color=WHITE)
    add_text(sl, name, bx+0.15, by+0.28, 3.7, 0.35, size=14, bold=True, color=WHITE)
    divider(sl, bx+0.15, by+0.56, 3.7, RGBColor(0xE5,0xE7,0xEB))
    add_text(sl, desc, bx+0.15, by+0.65, 3.75, 1.45, size=11, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — KEY BENEFITS & VALUE PROPOSITION
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "KEY BENEFITS", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Value Delivered by the Platform", 0.45, 0.42, 11, 0.7,
         size=30, bold=True, color=WHITE)

benefits = [
    ("🔐","Secure Dual Access",
     "Separate OTP flows for guests (Gmail) and employees (SMS). "
     "No password sharing, no unauthorised entry.",NAVY),
    ("📊","Real-Time Visibility",
     "Live project progress bars, active site counts, and employee "
     "status updated instantly across the platform.",BLUE),
    ("🏗️","Org Transparency",
     "Visual org chart from CEO to department heads, with clear "
     "reporting lines and team details for every leader.",GOLD),
    ("👥","Staff Directory",
     "Searchable employee list with department, ID, and Active / On Leave "
     "status — built for HR and operations teams.",SUCCESS),
    ("📅","Event Awareness",
     "Featured events, filter by category, RSVP tracking, and "
     "real-time notifications keep everyone aligned.",PURPLE),
    ("📱","Cross-Platform",
     "One React Native codebase runs natively on iOS and Android "
     "with identical UX and performance.",TEAL),
    ("📋","Projects at a Glance",
     "Project portfolio with progress tracking, site count, workforce "
     "and milestone dates — visible to management in seconds.",ERROR),
    ("🌐","Live Demo Ready",
     "Full interactive HTML prototype hosted on GitHub Pages for "
     "instant stakeholder access — no install required.",RGBColor(0xF5,0x9E,0x0B)),
]

for i,(icon,title,desc,c) in enumerate(benefits):
    bx = 0.45 + (i%4)*3.21
    by = 1.65 + (i//4)*2.7
    card_rect(sl, bx, by, 3.05, 2.45, WHITE)
    add_rect(sl, bx, by, 3.05, 0.1, c)
    add_text(sl, icon, bx+0.15, by+0.15, 0.55, 0.55, size=22, color=WHITE)
    add_text(sl, title, bx+0.15, by+0.68, 2.75, 0.42, size=13, bold=True,
             color=NAVY_DARK)
    add_text(sl, desc,  bx+0.15, by+1.08, 2.75, 1.1, size=10.5, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — GUEST vs EMPLOYEE ACCESS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "ACCESS COMPARISON", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "Guest vs Employee — What Each User Sees", 0.45, 0.42, 12, 0.7,
         size=28, bold=True, color=WHITE)

# GUEST card (left)
card_rect(sl, 0.45, 1.65, 6.2, 5.55, WHITE)
add_rect(sl, 0.45, 1.65, 6.2, 1.05, GOLD)
av_g = sl.shapes.add_shape(9, Inches(0.7), Inches(1.83), Inches(0.7), Inches(0.7))
av_g.fill.solid(); av_g.fill.fore_color.rgb = WHITE; av_g.line.fill.background()
add_text(sl, "G", 0.7, 1.88, 0.7, 0.6, size=22, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER)
add_text(sl, "GUEST USER", 1.55, 1.78, 4.5, 0.35, size=11, bold=True, color=NAVY_DARK)
add_text(sl, "View-Only Access", 1.55, 2.1,  4.5, 0.4, size=18, bold=True, color=WHITE)
add_text(sl, "Authenticated via Gmail OTP  ·  No account required",
         1.55, 2.42, 5.0, 0.28, size=10, color=WHITE)

add_text(sl, "WHAT GUESTS CAN VIEW", 0.65, 2.92, 5.5, 0.3, size=9, bold=True, color=GOLD)

guest_features = [
    "Company website content — About, Vision, Mission",
    "Services portfolio (Hard, Soft, Managed)",
    "Public projects gallery & client showcase",
    "Management team org chart",
    "Department heads overview",
    "Published news & company events",
    "Contact information & office locations",
    "ISO certifications & company milestones",
]
for i, text in enumerate(guest_features):
    by = 3.3 + i * 0.47
    add_text(sl, "✓", 0.7, by, 0.3, 0.35, size=14, bold=True, color=SUCCESS)
    add_text(sl, text, 1.0, by, 5.4, 0.35, size=11, color=TEXT_SEC)

# EMPLOYEE card (right)
card_rect(sl, 6.85, 1.65, 6.1, 5.55, WHITE)
add_rect(sl, 6.85, 1.65, 6.1, 1.05, NAVY)
av_e = sl.shapes.add_shape(9, Inches(7.1), Inches(1.83), Inches(0.7), Inches(0.7))
av_e.fill.solid(); av_e.fill.fore_color.rgb = GOLD; av_e.line.fill.background()
add_text(sl, "E", 7.1, 1.88, 0.7, 0.6, size=22, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER)
add_text(sl, "EMPLOYEE USER", 7.95, 1.78, 4.5, 0.35, size=11, bold=True, color=GOLD_LIGHT)
add_text(sl, "Full Self-Service Portal", 7.95, 2.1,  4.5, 0.4, size=18, bold=True, color=WHITE)
add_text(sl, "Authenticated via Employee ID + SMS OTP",
         7.95, 2.42, 5.0, 0.28, size=10, color=RGBColor(0xCC,0xDD,0xEE))

add_text(sl, "EVERYTHING GUESTS SEE  +  PERSONAL HR PORTAL", 7.05, 2.92, 5.7, 0.3,
         size=9, bold=True, color=GOLD)

emp_features = [
    "Personal profile · Joining date · Designation",
    "Leave balance · Apply leave · Approval status",
    "Pay slips — current & historical downloads",
    "Signed offer letter & employment contract",
    "Request salary / employment certificates",
    "Submit complaints & grievances directly to HR",
    "Access HR forms — claims, transfers, training",
    "Personalised notifications & work updates",
]
for i, text in enumerate(emp_features):
    by = 3.3 + i * 0.47
    add_text(sl, "★", 7.05, by, 0.35, 0.35, size=13, bold=True, color=GOLD)
    add_text(sl, text, 7.45, by, 5.45, 0.35, size=11, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — EMPLOYEE SELF-SERVICE PORTAL
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "EMPLOYEE SELF-SERVICE", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "HR Portal Features for Employees", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

ess_features = [
    ("Joining & Tenure",
     "View confirmed joining date, probation status, "
     "years of service, and upcoming work anniversaries.", NAVY),
    ("Leave Balance",
     "Live balance: Annual · Sick · Casual · Emergency. "
     "View entitlement, used, pending and remaining days.", SUCCESS),
    ("Digital Pay Slips",
     "Download current and historical pay slips as PDF. "
     "Breakdown of basic, allowances, deductions & net pay.", GOLD),
    ("Offer Letter & Contract",
     "Securely view your signed offer letter, employment "
     "contract, and any amendments — all in one place.", PURPLE),
    ("Salary Certificate",
     "Request salary / employment certificates for visa, "
     "bank loans, or rental — issued digitally by HR.", BLUE),
    ("Complaints & Grievances",
     "Submit confidential complaints to HR with categories, "
     "attach evidence, and track resolution status.", ERROR),
    ("HR Forms Library",
     "Access all official forms: leave, reimbursement, training, "
     "transfer requests, family addition, NOC, and more.", TEAL),
    ("Leave History",
     "Full audit trail of past leaves — dates, type, approval "
     "chain, and remarks. Filter by year or category.", RGBColor(0xF5,0x9E,0x0B)),
    ("HR Notifications",
     "Push alerts for approvals, payslip availability, "
     "policy updates, and HR announcements.", NAVY_LIGHT),
]

for i, (title, desc, c) in enumerate(ess_features):
    bx = 0.45 + (i % 3) * 4.3
    by = 1.65 + (i // 3) * 1.85
    card_rect(sl, bx, by, 4.05, 1.7, WHITE)
    add_rect(sl, bx, by, 4.05, 0.1, c)
    ic = sl.shapes.add_shape(9, Inches(bx + 0.18), Inches(by + 0.22),
                              Inches(0.55), Inches(0.55))
    ic.fill.solid(); ic.fill.fore_color.rgb = c; ic.line.fill.background()
    add_text(sl, str(i + 1), bx + 0.18, by + 0.3, 0.55, 0.4, size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx + 0.85, by + 0.28, 3.1, 0.4, size=13, bold=True,
             color=NAVY_DARK)
    add_text(sl, desc,  bx + 0.18, by + 0.85, 3.8, 0.85, size=10.5, color=TEXT_SEC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — LEAVE APPROVAL FLOW
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 1.45, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

add_text(sl, "LEAVE APPROVAL FLOW", 0.45, 0.12, 8, 0.32, size=10, bold=True, color=GOLD)
add_text(sl, "From Request to Approval — End-to-End", 0.45, 0.42, 11, 0.7,
         size=28, bold=True, color=WHITE)

# Left-side: live leave balance card
card_rect(sl, 0.45, 1.65, 3.7, 5.55, WHITE)
add_rect(sl, 0.45, 1.65, 3.7, 0.5, NAVY)
add_text(sl, "MY LEAVE BALANCE", 0.6, 1.74, 3.5, 0.35, size=11, bold=True,
         color=GOLD_LIGHT)

bals = [
    ("Annual Leave", 18, 30, SUCCESS),
    ("Sick Leave",    4, 14, BLUE),
    ("Casual Leave",  2,  7, GOLD),
    ("Emergency",     1,  3, ERROR),
]
for i, (lbl, used, total, c) in enumerate(bals):
    by = 2.32 + i * 1.18
    add_text(sl, lbl, 0.65, by, 3.4, 0.3, size=11, bold=True, color=NAVY_DARK)
    rem = total - used
    add_text(sl, str(rem) + " days remaining", 0.65, by + 0.3, 3.4, 0.3,
             size=10, color=TEXT_SEC)
    add_rect(sl, 0.65, by + 0.65, 3.4, 0.18, RGBColor(0xE5, 0xE7, 0xEB))
    fill_w = 3.4 * used / total
    add_rect(sl, 0.65, by + 0.65, fill_w, 0.18, c)
    add_text(sl, "Used " + str(used) + " of " + str(total),
             0.65, by + 0.86, 3.4, 0.25, size=9, color=TEXT_LIGHT,
             align=PP_ALIGN.RIGHT)

# Apply button
add_rect(sl, 0.65, 6.7, 3.3, 0.4, GOLD)
add_text(sl, "+ Apply for Leave", 0.65, 6.74, 3.3, 0.32, size=11, bold=True,
         color=NAVY_DARK, align=PP_ALIGN.CENTER)

# Right-side: approval flow steps
flow_steps = [
    ("01", "Employee Submits Request",
     "Choose leave type, date range, attach reason. App auto-calculates "
     "days deducted and warns of overlapping balance.", NAVY, SUCCESS, "Submitted"),
    ("02", "Direct Manager Review",
     "Real-time push notification to line manager. Manager approves, "
     "rejects, or requests clarification — all in-app.", NAVY_LIGHT, SUCCESS, "Approved"),
    ("03", "Department Head Approval",
     "For leaves > 5 days, escalated to Department Head with full "
     "team-coverage context shown.", GOLD, SUCCESS, "Approved"),
    ("04", "HR Final Sign-Off",
     "HR validates balance, updates leave records, and triggers "
     "calendar block. Confirmation sent to employee.", NAVY, GOLD, "Pending"),
    ("05", "Status & Audit Trail",
     "Employee sees live status: Pending to Approved to Confirmed. "
     "Full audit trail preserved for compliance.", PURPLE,
     RGBColor(0xE5,0xE7,0xEB), "Awaiting"),
]

for i, (num, title, desc, nc, sc, slbl) in enumerate(flow_steps):
    by = 1.65 + i * 1.13
    card_rect(sl, 4.35, by, 8.6, 0.98, WHITE)
    ic = sl.shapes.add_shape(9, Inches(4.45), Inches(by + 0.18),
                              Inches(0.62), Inches(0.62))
    ic.fill.solid(); ic.fill.fore_color.rgb = nc; ic.line.fill.background()
    add_text(sl, num, 4.45, by + 0.28, 0.62, 0.45, size=14, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, title, 5.2, by + 0.08, 5.5, 0.38, size=13, bold=True,
             color=NAVY_DARK)
    add_text(sl, desc,  5.2, by + 0.46, 7.5, 0.5, size=10, color=TEXT_SEC)
    badge = sl.shapes.add_shape(5, Inches(11.45), Inches(by + 0.32),
                                 Inches(1.4), Inches(0.34))
    badge.fill.solid(); badge.fill.fore_color.rgb = sc
    badge.line.fill.background()
    sc_text_col = WHITE if sc != RGBColor(0xE5,0xE7,0xEB) else TEXT_SEC
    add_text(sl, slbl, 11.45, by + 0.34, 1.4, 0.32, size=10, bold=True,
             color=sc_text_col, align=PP_ALIGN.CENTER)
    if i < 4:
        add_rect(sl, 4.7, by + 0.99, 0.1, 0.14, GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — NEXT STEPS & CLOSING
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, GOLD)

# Decorative circles
c1 = sl.shapes.add_shape(9, Inches(9.5), Inches(-1.2), Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x22,0x44,0x7A)
c1.line.fill.background()
c2 = sl.shapes.add_shape(9, Inches(-0.6), Inches(5.5), Inches(3), Inches(3))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xC9,0xA8,0x4C)
c2.line.fill.background()

add_text(sl, "NEXT STEPS", 0.45, 0.85, 5, 0.32, size=10, bold=True, color=GOLD)
add_rect(sl, 0.45, 1.18, 1.8, 0.055, GOLD)
add_text(sl, "Roadmap & Call to Action", 0.45, 1.28, 10, 0.65, size=32,
         bold=True, color=WHITE)

steps_ns = [
    ("Q2 2026","Stakeholder Review & Sign-off",
     "Present live demo to all department heads. Gather feedback and document change requests."),
    ("Q3 2026","Backend Integration",
     "Connect app to CAFM/CMMS APIs. Enable live work-order creation, push notifications, and data sync."),
    ("Q3 2026","UAT & Employee Onboarding",
     "Pilot with 50 employees. Collect usability metrics, conduct training sessions, refine UX."),
    ("Q4 2026","Full Production Launch",
     "Go-live for all 3,500+ staff. App Store (iOS) and Google Play (Android) submission."),
]

for i,(quarter,title,desc) in enumerate(steps_ns):
    by = 2.15 + i * 1.15
    add_rect(sl, 0.45, by, 1.45, 0.95, RGBColor(0x1F,0x42,0x75))
    add_text(sl, quarter, 0.45, by+0.15, 1.45, 0.55, size=12, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.95, by+0.4, 0.02, 0.12, GOLD)
    add_rect(sl, 2.0, by+0.43, 0.25, 0.06, GOLD)
    add_text(sl, title, 2.35, by+0.06, 8.5, 0.38, size=14, bold=True, color=WHITE)
    add_text(sl, desc,  2.35, by+0.44, 8.5, 0.42, size=11, color=RGBColor(0xAA,0xBB,0xCC))

# CTA
add_rect(sl, 0.45, 6.6 - 0.3, 5.8, 0.6, RGBColor(0x1F,0x42,0x75))
add_text(sl, "🌐  Live Demo Available Now",
         0.6, 6.32, 5.5, 0.52, size=13, bold=True, color=GOLD_LIGHT)

add_rect(sl, 6.65, 6.3, 6.3, 0.6, GOLD)
add_text(sl, "mebikramshah-design.github.io/Mobille-App-Final-",
         6.75, 6.33, 6.1, 0.52, size=12, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER)

# Thank you
add_text(sl, "Thank You", 0.45, 5.6, 12, 0.72, size=48, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Darwish Interserve Facility Management  ·  Qatar  ·  2026",
         0.45, 6.08, 12.4, 0.35, size=12, color=RGBColor(0x88,0x99,0xAA),
         align=PP_ALIGN.CENTER)

# Footer bar
add_rect(sl, 0, 7.15, 13.33, 0.35, RGBColor(0x0D,0x1C,0x35))
add_text(sl, "CONFIDENTIAL — FOR MANAGEMENT USE ONLY", 0, 7.18, 13.33, 0.3,
         size=9, color=RGBColor(0x55,0x66,0x77), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = "/home/user/Mobille-App-Final-/DarwishInterserve_Executive_Presentation.pptx"
prs.save(out)
print("Saved:", out)
